"""Build a single PPTX slide summarizing the Joint Account KYC Basic ID upload flow for the Alloy team."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Colors
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
BLUE_BG = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_BORDER = RGBColor(0x25, 0x63, 0xEB)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
AMBER_BORDER = RGBColor(0xD9, 0x77, 0x06)
AMBER_TEXT = RGBColor(0x78, 0x35, 0x0F)
GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_BORDER = RGBColor(0x16, 0xA3, 0x4A)
GREEN_TEXT = RGBColor(0x14, 0x53, 0x2D)
PURPLE_BG = RGBColor(0xED, 0xE9, 0xFE)
PURPLE_BORDER = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_TEXT = RGBColor(0x4C, 0x1D, 0x95)
GRAY_TEXT = RGBColor(0x37, 0x41, 0x51)
LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color, width_pt=1.25):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_text(shape, text, size=10, bold=False, color=GRAY_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"


def add_box(slide, left, top, width, height, fill, border, text, size=10, bold=False, text_color=GRAY_TEXT, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_fill(shp, fill)
    set_line(shp, border, 1.25)
    if radius:
        # Slight corner rounding
        shp.adjustments[0] = 0.1
    add_text(shp, text, size=size, bold=bold, color=text_color)
    return shp


def add_connector(slide, x1, y1, x2, y2, color=GRAY_TEXT, width_pt=1.5, dashed=False, label=None):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    # Arrow end
    ln = conn.line._get_or_add_ln()
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    if dashed:
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    if label:
        # Small textbox label near midpoint
        mx = (x1 + x2) // 2
        my = (y1 + y2) // 2
        tb = slide.shapes.add_textbox(mx - Inches(0.9), my - Inches(0.12), Inches(1.8), Inches(0.25))
        tf = tb.text_frame
        tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
        tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(8); r.font.italic = True; r.font.color.rgb = color; r.font.name = "Calibri"
    return conn


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    title_tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.5))
    tf = title_tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Joint Account Flow — KYC Basic — ID Upload Sequencing"
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"

    sub_tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.35))
    tf = sub_tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "For Alloy journey config: where the primary user and each co-holder upload their ID in a KYC Basic session"
    r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = GRAY_TEXT; r.font.name = "Calibri"

    # === PRIMARY USER lane (left) ===
    lane_left = Inches(0.4)
    lane_top = Inches(1.3)
    lane_w = Inches(6.2)
    lane_h = Inches(5.3)

    # Lane background
    primary_lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lane_left, lane_top, lane_w, lane_h)
    set_fill(primary_lane, RGBColor(0xF5, 0xF9, 0xFF))
    set_line(primary_lane, BLUE_BORDER, 1.5)
    primary_lane.adjustments[0] = 0.03
    # Lane header
    hdr = slide.shapes.add_textbox(lane_left, lane_top + Inches(0.05), lane_w, Inches(0.35))
    tf = hdr.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "PRIMARY USER — inside embedded component"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"

    # Primary steps boxes
    box_w = Inches(5.7); box_h = Inches(0.5)
    bx = lane_left + Inches(0.25)
    by = lane_top + Inches(0.5)
    gap = Inches(0.12)

    steps = [
        ("Screen 5 · Step 1 — Identity info (name, DOB, SSN)", BLUE_BG, BLUE_BORDER, NAVY, False),
        ("Screen 6 · Step 2 — Address", BLUE_BG, BLUE_BORDER, NAVY, False),
        ("Screen 7 · Step 3 — Proof of address upload", BLUE_BG, BLUE_BORDER, NAVY, False),
        ("Screen 8a · Step 4 — Issuing country + document type", BLUE_BG, BLUE_BORDER, NAVY, False),
        ("Screen 8b · Step 4 — PRIMARY ID UPLOAD (image(s))", GREEN_BG, GREEN_BORDER, GREEN_TEXT, True),
        ("Screen 9 / 9a · Step 5 — Enter co-holder bio data (NO ID upload here)", BLUE_BG, BLUE_BORDER, NAVY, False),
        ("Screen 10 — Distribute co-holder verification links", BLUE_BG, BLUE_BORDER, NAVY, False),
    ]
    primary_shapes = []
    for (txt, fill, border, tcolor, bold) in steps:
        shp = add_box(slide, bx, by, box_w, box_h, fill, border, txt, size=10.5, bold=bold, text_color=tcolor)
        primary_shapes.append(shp)
        by += box_h + gap

    # Arrows between primary steps
    for i in range(len(primary_shapes) - 1):
        s1 = primary_shapes[i]; s2 = primary_shapes[i+1]
        x = s1.left + s1.width // 2
        add_connector(slide, x, s1.top + s1.height, x, s2.top, color=BLUE_BORDER, width_pt=1.25)

    # === CO-HOLDER lane (right) ===
    co_left = Inches(7.0)
    co_top = lane_top
    co_w = Inches(5.9)
    co_h = lane_h

    co_lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, co_left, co_top, co_w, co_h)
    set_fill(co_lane, RGBColor(0xFF, 0xFB, 0xEB))
    set_line(co_lane, AMBER_BORDER, 1.5)
    co_lane.adjustments[0] = 0.03

    hdr2 = slide.shapes.add_textbox(co_left, co_top + Inches(0.05), co_w, Inches(0.55))
    tf = hdr2.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = 0; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "EACH CO-HOLDER — standalone token-authenticated link"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = AMBER_TEXT; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = "(repeats independently for each co-holder, 1–4 times)"
    r2.font.size = Pt(9); r2.font.italic = True; r2.font.color.rgb = AMBER_TEXT; r2.font.name = "Calibri"

    # Delivery options (3 small boxes side by side)
    deliv_y = co_top + Inches(0.7)
    deliv_w = Inches(1.75); deliv_h = Inches(0.75)
    deliv_x = co_left + Inches(0.15)

    deliveries = [
        ("Send to email\n(primary path)", AMBER_BG, AMBER_BORDER, AMBER_TEXT),
        ("Copy link", AMBER_BG, AMBER_BORDER, AMBER_TEXT),
        ("Verify now\n(edge case)", LIGHT_GRAY, AMBER_BORDER, GRAY_TEXT),
    ]
    deliv_shapes = []
    for (txt, fill, border, tc) in deliveries:
        shp = add_box(slide, deliv_x, deliv_y, deliv_w, deliv_h, fill, border, txt, size=9.5, bold=True, text_color=tc)
        deliv_shapes.append(shp)
        deliv_x += deliv_w + Inches(0.1)

    # Co-holder journey boxes
    co_bx = co_left + Inches(0.25)
    co_by = deliv_y + deliv_h + Inches(0.3)
    co_box_w = Inches(5.4)

    co_steps = [
        ("Standalone page at {client}.interro.co\n(token-authenticated, host-branded)", AMBER_BG, AMBER_BORDER, AMBER_TEXT, 0.55, False),
        ("Co-holder Screen 8a — Issuing country + document type", AMBER_BG, AMBER_BORDER, AMBER_TEXT, 0.5, False),
        ("Co-holder Screen 8b — CO-HOLDER ID UPLOAD (image(s))", GREEN_BG, GREEN_BORDER, GREEN_TEXT, 0.5, True),
        ("Co-holder submission → webhook to host + Interro", AMBER_BG, AMBER_BORDER, AMBER_TEXT, 0.5, False),
    ]
    co_shapes = []
    for (txt, fill, border, tc, h_in, bold) in co_steps:
        shp = add_box(slide, co_bx, co_by, co_box_w, Inches(h_in), fill, border, txt, size=10, bold=bold, text_color=tc)
        co_shapes.append(shp)
        co_by += Inches(h_in) + gap

    # Arrows in co-holder lane
    # deliveries converge to first co step
    first_co = co_shapes[0]
    target_x = first_co.left + first_co.width // 2
    target_y = first_co.top
    for ds in deliv_shapes:
        sx = ds.left + ds.width // 2
        sy = ds.top + ds.height
        add_connector(slide, sx, sy, target_x, target_y, color=AMBER_BORDER, width_pt=1.25)

    for i in range(len(co_shapes) - 1):
        s1 = co_shapes[i]; s2 = co_shapes[i+1]
        x = s1.left + s1.width // 2
        add_connector(slide, x, s1.top + s1.height, x, s2.top, color=AMBER_BORDER, width_pt=1.25)

    # Cross-lane arrow: Screen 10 (primary) → delivery options (co lane)
    s10 = primary_shapes[-1]
    # from right edge of s10 to left edge of middle delivery box
    src_x = s10.left + s10.width
    src_y = s10.top + s10.height // 2
    mid_deliv = deliv_shapes[1]
    dst_x = mid_deliv.left + mid_deliv.width // 2
    dst_y = mid_deliv.top
    add_connector(slide, src_x, src_y, dst_x, dst_y, color=GRAY_TEXT, width_pt=1.5)

    # Alloy callouts (dashed) — primary ID upload → Alloy applicant #1
    s8b = primary_shapes[4]
    alloy_p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(6.55), s8b.top - Inches(0.05),
                                     Inches(0.5), Inches(0.6))
    set_fill(alloy_p, PURPLE_BG); set_line(alloy_p, PURPLE_BORDER, 1.25)
    alloy_p.adjustments[0] = 0.15
    add_text(alloy_p, "Alloy\napp #1", size=8, bold=True, color=PURPLE_TEXT)

    # Co-holder upload → Alloy per-person
    ch_upload = co_shapes[2]
    alloy_c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(12.75), ch_upload.top - Inches(0.05),
                                     Inches(0.5), Inches(0.6))
    # Place it inside the lane instead to avoid overflow
    alloy_c.left = Inches(12.55)
    set_fill(alloy_c, PURPLE_BG); set_line(alloy_c, PURPLE_BORDER, 1.25)
    alloy_c.adjustments[0] = 0.15
    add_text(alloy_c, "Alloy\n1 app per\nco-holder", size=7.5, bold=True, color=PURPLE_TEXT)

    # Dashed arrows from upload boxes to Alloy callouts
    add_connector(slide,
                  s8b.left + s8b.width, s8b.top + s8b.height // 2,
                  alloy_p.left, alloy_p.top + alloy_p.height // 2,
                  color=PURPLE_BORDER, width_pt=1.25, dashed=True)
    add_connector(slide,
                  ch_upload.left + ch_upload.width, ch_upload.top + ch_upload.height // 2,
                  alloy_c.left, alloy_c.top + alloy_c.height // 2,
                  color=PURPLE_BORDER, width_pt=1.25, dashed=True)

    # Footer key takeaways
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.65))
    tf = footer.text_frame; tf.margin_left = 0; tf.margin_top = 0; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Alloy config implications: "
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = ("one journey, one application per holder (primary + each co-holder), "
               "keyed by person index and session ID. Applications arrive asynchronously "
               "(minutes to days apart). KYC Basic does NOT use Alloy SDK live capture — "
               "journey must accept pre-captured document images + biographical data.")
    r2.font.size = Pt(10); r2.font.color.rgb = GRAY_TEXT; r2.font.name = "Calibri"

    out = r"C:\Users\MarcoCesaratto\Desktop\Embeddable-UI\designs\joint-account-kyc-basic-id-upload-flow.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
